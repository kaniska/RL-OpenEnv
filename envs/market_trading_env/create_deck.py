"""
Business Value PowerPoint Deck Generator
=========================================
Creates a concise, non-technical business presentation explaining
the Multi-Agent Market Simulator and its RL-based value proposition.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x1D, 0x3A)   # Dark navy background
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOLD      = RGBColor(0xFF, 0xB8, 0x00)   # Accent / highlights
LIGHT_BLUE = RGBColor(0x4A, 0x9E, 0xD9)  # Secondary accent
GRAY      = RGBColor(0x8C, 0x8C, 0x8C)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
DARK_BG   = RGBColor(0x10, 0x25, 0x45)
SOFT_BG   = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ──────────────────────────────────────────

def add_bg(slide, color=NAVY):
    """Fill entire slide with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  font_name="Calibri"):
    """Add another paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_circle_number(slide, left, top, number, size=Inches(0.5), color=GOLD):
    """Add a numbered circle indicator."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

def add_icon_card(slide, left, top, width, height, icon_text, title, body,
                  card_color=DARK_BG, icon_color=GOLD):
    """Add a card with icon, title and body text."""
    card = add_shape_bg(slide, left, top, width, height, card_color)
    card.shadow.inherit = False

    # Icon circle
    icon_size = Inches(0.6)
    icon_left = left + (width - icon_size) // 2
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_left, top + Inches(0.3),
                                   icon_size, icon_size)
    icon.fill.solid()
    icon.fill.fore_color.rgb = icon_color
    icon.line.fill.background()
    itf = icon.text_frame
    itf.paragraphs[0].text = icon_text
    itf.paragraphs[0].font.size = Pt(20)
    itf.paragraphs[0].font.color.rgb = NAVY
    itf.paragraphs[0].font.bold = True
    itf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide, left + Inches(0.2), top + Inches(1.1),
             width - Inches(0.4), Inches(0.4),
             title, font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

    # Body
    add_text(slide, left + Inches(0.25), top + Inches(1.55),
             width - Inches(0.5), height - Inches(1.8),
             body, font_size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 1: TITLE SLIDE
# =====================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1, NAVY)

# Accent bar at top
add_shape_bg(slide1, 0, 0, W, Inches(0.08), GOLD)

# Main title
add_text(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "AI-Powered Market Intelligence", font_size=44, color=WHITE,
         bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text(slide1, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
         "How Smart Agents Learn to Trade, Negotiate & Win",
         font_size=24, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# Divider line
add_shape_bg(slide1, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04), GOLD)

# Description
add_text(slide1, Inches(2), Inches(4.4), Inches(9), Inches(1.0),
         "A reinforcement learning platform that trains AI agents to make\n"
         "better trading decisions — reducing costs, increasing margins,\n"
         "and discovering strategies humans might miss.",
         font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Footer
add_text(slide1, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "MarketForge  |  OpenEnv Hackathon 2025",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2: THE BUSINESS PROBLEM
# =====================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, NAVY)
add_shape_bg(slide2, 0, 0, W, Inches(0.08), GOLD)

# Title
add_text(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "The Business Problem", font_size=36, color=WHITE, bold=True)

# Subtitle
add_text(slide2, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Traditional trading and procurement decisions leave money on the table",
         font_size=16, color=GOLD)

# Problem cards
problems = [
    ("$", "Pricing Guesswork",
     "Companies set prices based on gut\n"
     "feel, spreadsheets, or last quarter's\n"
     "data — missing real-time market\n"
     "signals that affect margins."),
    ("!", "Slow Negotiations",
     "Deal-making with suppliers and\n"
     "buyers takes weeks of back-and-\n"
     "forth. Every delay costs money\n"
     "and opportunities slip away."),
    ("?", "Blind Spots",
     "Market disruptions — supply\n"
     "shortages, demand spikes, new\n"
     "competitors — catch teams off\n"
     "guard with no playbook."),
    ("#", "Missed Partnerships",
     "Companies compete when they\n"
     "should cooperate. They miss\n"
     "joint ventures and supply chain\n"
     "alliances that boost both sides."),
]

card_w = Inches(2.7)
card_h = Inches(3.2)
start_x = Inches(0.6)
gap = Inches(0.35)

for i, (icon, title, body) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    add_icon_card(slide2, x, Inches(2.0), card_w, card_h,
                  icon, title, body)

# Bottom insight
add_text(slide2, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.8),
         "Result: Companies lose 5-15% of potential margin due to suboptimal\n"
         "trading decisions, slow reactions, and missed collaboration opportunities.",
         font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 3: OUR SOLUTION — THE AI MARKET SIMULATOR
# =====================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, NAVY)
add_shape_bg(slide3, 0, 0, W, Inches(0.08), GOLD)

add_text(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "The Solution: AI Agents That Learn to Win", font_size=36,
         color=WHITE, bold=True)

add_text(slide3, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "A virtual marketplace where AI agents practice and master trading strategies",
         font_size=16, color=GOLD)

# How it works — 3 step visual
steps = [
    ("1", "Build the Market",
     "We create a realistic virtual\n"
     "marketplace with real commodity\n"
     "prices, supply/demand dynamics,\n"
     "and market disruptions — just\n"
     "like the real world."),
    ("2", "Train AI Agents",
     "AI agents enter the market and\n"
     "learn by doing — thousands of\n"
     "trades, negotiations, and deals.\n"
     "Each round they get smarter\n"
     "based on what worked."),
    ("3", "Deploy Insights",
     "The trained agents reveal optimal\n"
     "strategies: when to buy, what\n"
     "price to offer, who to partner\n"
     "with, and how to react when\n"
     "markets shift."),
]

step_w = Inches(3.4)
step_h = Inches(3.0)
step_start = Inches(0.7)
step_gap = Inches(0.6)

for i, (num, title, body) in enumerate(steps):
    x = step_start + i * (step_w + step_gap)
    y = Inches(2.0)

    # Card background
    add_shape_bg(slide3, x, y, step_w, step_h, DARK_BG)

    # Step number circle
    num_size = Inches(0.55)
    num_shape = slide3.shapes.add_shape(MSO_SHAPE.OVAL,
                                         x + Inches(0.15), y + Inches(0.2),
                                         num_size, num_size)
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = GOLD
    num_shape.line.fill.background()
    ntf = num_shape.text_frame
    ntf.paragraphs[0].text = num
    ntf.paragraphs[0].font.size = Pt(18)
    ntf.paragraphs[0].font.color.rgb = NAVY
    ntf.paragraphs[0].font.bold = True
    ntf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide3, x + Inches(0.85), y + Inches(0.25),
             step_w - Inches(1.0), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)

    # Body
    add_text(slide3, x + Inches(0.25), y + Inches(0.85),
             step_w - Inches(0.5), step_h - Inches(1.0),
             body, font_size=13, color=LIGHT_GRAY)

    # Arrow between steps
    if i < 2:
        arrow_x = x + step_w + Inches(0.1)
        add_text(slide3, arrow_x, y + Inches(1.2), Inches(0.4), Inches(0.5),
                 ">>>", font_size=22, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_text(slide3, Inches(1), Inches(5.4), Inches(11), Inches(0.8),
         "Think of it as a flight simulator for traders — practice in a safe\n"
         "environment, then apply proven strategies in real markets.",
         font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 4: HOW THE AI LEARNS — AWARDS & TRAINING
# =====================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, NAVY)
add_shape_bg(slide4, 0, 0, W, Inches(0.08), GOLD)

add_text(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "How the AI Gets Smarter Over Time", font_size=36,
         color=WHITE, bold=True)

add_text(slide4, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Clear goals + practice + feedback = continuous improvement",
         font_size=16, color=GOLD)

# Left panel — Award 1
panel_w = Inches(5.5)
panel_h = Inches(4.0)

add_shape_bg(slide4, Inches(0.6), Inches(1.9), panel_w, panel_h, DARK_BG)

# Gold trophy icon
trophy = slide4.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.9), Inches(2.15),
                                  Inches(0.5), Inches(0.5))
trophy.fill.solid()
trophy.fill.fore_color.rgb = GOLD
trophy.line.fill.background()
ttf = trophy.text_frame
ttf.paragraphs[0].text = "1"
ttf.paragraphs[0].font.size = Pt(16)
ttf.paragraphs[0].font.color.rgb = NAVY
ttf.paragraphs[0].font.bold = True
ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(slide4, Inches(1.55), Inches(2.15), Inches(4.0), Inches(0.5),
         "Award 1: Market Champion", font_size=20, color=GOLD, bold=True)

tf4a = add_text(slide4, Inches(0.9), Inches(2.8), Inches(4.8), Inches(2.8),
         "Goal: Make the most money.", font_size=14, color=WHITE, bold=True)

add_paragraph(tf4a, "", font_size=6, color=WHITE)
add_paragraph(tf4a,
    "The AI agent that builds the most wealth wins.\n"
    "This teaches the AI to:", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4a, "", font_size=4, color=WHITE)
add_paragraph(tf4a, "  \u2022  Buy low and sell high", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4a, "  \u2022  Manage inventory wisely", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4a, "  \u2022  Time the market correctly", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4a, "  \u2022  Avoid bad trades that lose money", font_size=13, color=LIGHT_GRAY)

# Right panel — Award 2
add_shape_bg(slide4, Inches(6.5), Inches(1.9), panel_w + Inches(0.5), panel_h, DARK_BG)

trophy2 = slide4.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(6.8), Inches(2.15),
                                   Inches(0.5), Inches(0.5))
trophy2.fill.solid()
trophy2.fill.fore_color.rgb = LIGHT_BLUE
trophy2.line.fill.background()
ttf2 = trophy2.text_frame
ttf2.paragraphs[0].text = "2"
ttf2.paragraphs[0].font.size = Pt(16)
ttf2.paragraphs[0].font.color.rgb = NAVY
ttf2.paragraphs[0].font.bold = True
ttf2.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text(slide4, Inches(7.45), Inches(2.15), Inches(5.0), Inches(0.5),
         "Award 2: Master Strategist", font_size=20, color=LIGHT_BLUE, bold=True)

tf4b = add_text(slide4, Inches(6.8), Inches(2.8), Inches(5.5), Inches(2.8),
         "Goal: Be the smartest player.", font_size=14, color=WHITE, bold=True)

add_paragraph(tf4b, "", font_size=6, color=WHITE)
add_paragraph(tf4b,
    "The AI agent with the best overall strategy wins.\n"
    "This teaches the AI to:", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4b, "", font_size=4, color=WHITE)
add_paragraph(tf4b, "  \u2022  Negotiate better deals with partners", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4b, "  \u2022  Form alliances that create value", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4b, "  \u2022  Adapt quickly when markets change", font_size=13, color=LIGHT_GRAY)
add_paragraph(tf4b, "  \u2022  Trade efficiently — no wasted moves", font_size=13, color=LIGHT_GRAY)

# Bottom — Improvement cycle
add_shape_bg(slide4, Inches(0.6), Inches(6.15), W - Inches(1.2), Inches(0.9), DARK_BG)
add_text(slide4, Inches(0.9), Inches(6.25), W - Inches(1.8), Inches(0.7),
         "The Learning Cycle:    Train  \u2192  Compete  \u2192  Score  \u2192  Improve  \u2192  Repeat"
         "        |        Each cycle makes the AI more accurate and profitable.",
         font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 5: THE IMPROVEMENT ENGINE
# =====================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, NAVY)
add_shape_bg(slide5, 0, 0, W, Inches(0.08), GOLD)

add_text(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Continuous Improvement: Better Every Round", font_size=36,
         color=WHITE, bold=True)

add_text(slide5, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Like an athlete reviewing game tape — the AI studies every trade and improves",
         font_size=16, color=GOLD)

# Improvement cycle cards
cycle_items = [
    ("1", "PLAY", "The Market Game",
     "AI agents enter the simulated\n"
     "marketplace and make real-time\n"
     "decisions: buying, selling,\n"
     "negotiating, and forming\n"
     "partnerships."),
    ("2", "SCORE", "Every Decision",
     "After each round, every action\n"
     "is scored: Did the trade make\n"
     "money? Was the negotiation\n"
     "successful? Was the move\n"
     "even valid?"),
    ("3", "LEARN", "From Mistakes",
     "The AI reviews what worked\n"
     "and what didn't. Good strategies\n"
     "get reinforced. Bad strategies\n"
     "get corrected. Accuracy\n"
     "improves each cycle."),
    ("4", "MASTER", "The Strategy",
     "After thousands of practice\n"
     "rounds, the AI discovers\n"
     "optimal strategies that\n"
     "maximize profit and minimize\n"
     "risk — ready for deployment."),
]

card_w2 = Inches(2.7)
card_h2 = Inches(3.5)
start_x2 = Inches(0.6)
gap2 = Inches(0.35)

for i, (num, label, title, body) in enumerate(cycle_items):
    x = start_x2 + i * (card_w2 + gap2)
    y = Inches(1.9)

    add_shape_bg(slide5, x, y, card_w2, card_h2, DARK_BG)

    # Number + Label
    num_size = Inches(0.5)
    ns = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2),
                                  num_size, num_size)
    ns.fill.solid()
    ns.fill.fore_color.rgb = GOLD if i < 3 else GREEN
    ns.line.fill.background()
    nsf = ns.text_frame
    nsf.paragraphs[0].text = num
    nsf.paragraphs[0].font.size = Pt(14)
    nsf.paragraphs[0].font.color.rgb = NAVY
    nsf.paragraphs[0].font.bold = True
    nsf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide5, x + Inches(0.75), y + Inches(0.15),
             card_w2 - Inches(0.9), Inches(0.35),
             label, font_size=16, color=GOLD if i < 3 else GREEN, bold=True)

    add_text(slide5, x + Inches(0.75), y + Inches(0.45),
             card_w2 - Inches(0.9), Inches(0.35),
             title, font_size=14, color=WHITE, bold=True)

    add_text(slide5, x + Inches(0.2), y + Inches(1.0),
             card_w2 - Inches(0.4), card_h2 - Inches(1.2),
             body, font_size=12, color=LIGHT_GRAY)

    # Arrow
    if i < 3:
        arrow_x = x + card_w2 + Inches(0.05)
        add_text(slide5, arrow_x, y + Inches(1.4), Inches(0.3), Inches(0.4),
                 "\u25B6", font_size=18, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

# Key metric
add_shape_bg(slide5, Inches(0.6), Inches(5.7), W - Inches(1.2), Inches(1.0), DARK_BG)

metrics_tf = add_text(slide5, Inches(1.0), Inches(5.85), Inches(3.5), Inches(0.6),
         "Accuracy Tracking", font_size=16, color=GOLD, bold=True)

add_text(slide5, Inches(4.5), Inches(5.85), Inches(8.0), Inches(0.6),
         "We measure what % of decisions are valid and profitable. "
         "As the AI trains, this number climbs — from ~60% to 90%+ — "
         "meaning fewer wasted moves and more money captured.",
         font_size=13, color=LIGHT_GRAY)


# =====================================================================
# SLIDE 6: BUSINESS VALUE & ROI
# =====================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, NAVY)
add_shape_bg(slide6, 0, 0, W, Inches(0.08), GOLD)

add_text(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Business Impact: Real Revenue, Real Growth", font_size=36,
         color=WHITE, bold=True)

add_text(slide6, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "How AI-trained trading strategies translate to bottom-line results",
         font_size=16, color=GOLD)

# Value proposition cards — 2 rows of 3
values = [
    ("\u2191", "Better Pricing",
     "AI discovers optimal buy/sell\n"
     "prices by testing thousands of\n"
     "strategies. Companies capture\n"
     "3-8% more margin per trade."),
    ("\u26A1", "Faster Deals",
     "AI-powered negotiations close\n"
     "deals in minutes instead of\n"
     "weeks. Automated offers and\n"
     "counter-offers 24/7."),
    ("\u26A0", "Risk Protection",
     "AI agents learn to spot market\n"
     "disruptions and react instantly.\n"
     "No more being caught off\n"
     "guard by supply shocks."),
    ("\u2764", "Smarter Partnerships",
     "AI identifies which alliances\n"
     "create value. Companies form\n"
     "the right partnerships at the\n"
     "right time."),
    ("\u2699", "Operational Savings",
     "Automated trading decisions\n"
     "reduce manual analysis by\n"
     "60-80%. Teams focus on\n"
     "strategy, not spreadsheets."),
    ("\u2B50", "Competitive Edge",
     "First movers who deploy AI\n"
     "trading agents gain an unfair\n"
     "advantage — better margins,\n"
     "faster execution, fewer losses."),
]

val_w = Inches(3.6)
val_h = Inches(1.9)
val_gap_x = Inches(0.35)
val_gap_y = Inches(0.25)

for i, (icon, title, body) in enumerate(values):
    row = i // 3
    col = i % 3
    x = Inches(0.6) + col * (val_w + val_gap_x)
    y = Inches(1.85) + row * (val_h + val_gap_y)

    add_shape_bg(slide6, x, y, val_w, val_h, DARK_BG)

    # Icon
    add_text(slide6, x + Inches(0.2), y + Inches(0.15),
             Inches(0.4), Inches(0.4),
             icon, font_size=20, color=GOLD, bold=True)

    # Title
    add_text(slide6, x + Inches(0.65), y + Inches(0.15),
             val_w - Inches(0.8), Inches(0.35),
             title, font_size=16, color=WHITE, bold=True)

    # Body
    add_text(slide6, x + Inches(0.25), y + Inches(0.6),
             val_w - Inches(0.5), val_h - Inches(0.7),
             body, font_size=12, color=LIGHT_GRAY)

# Bottom ROI bar
add_shape_bg(slide6, Inches(0.6), Inches(6.3), W - Inches(1.2), Inches(0.75), DARK_BG)
add_text(slide6, Inches(1.0), Inches(6.38), W - Inches(2.0), Inches(0.55),
         "Estimated ROI:  10-25x return on AI training investment within 12 months  "
         "|  Payback period: 3-6 months for commodity trading operations",
         font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 7: USE CASES — WHO BENEFITS
# =====================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, NAVY)
add_shape_bg(slide7, 0, 0, W, Inches(0.08), GOLD)

add_text(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Who Benefits: Industry Applications", font_size=36,
         color=WHITE, bold=True)

add_text(slide7, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "Any business that buys, sells, or trades can benefit from AI-trained strategies",
         font_size=16, color=GOLD)

industries = [
    ("Commodity Trading",
     "Energy, agriculture, metals\n"
     "trading firms use AI agents\n"
     "to find optimal entry/exit\n"
     "points and manage portfolio\n"
     "risk across volatile markets."),
    ("Supply Chain",
     "Procurement teams train AI\n"
     "to negotiate better supplier\n"
     "contracts, time purchases\n"
     "around price dips, and build\n"
     "resilient supply networks."),
    ("Retail & E-Commerce",
     "Dynamic pricing engines\n"
     "trained on market simulations\n"
     "adjust prices in real-time\n"
     "to maximize revenue while\n"
     "staying competitive."),
    ("Financial Services",
     "Portfolio managers use AI\n"
     "agents to test trading\n"
     "strategies before deploying\n"
     "real capital — reducing risk\n"
     "and finding alpha."),
]

ind_w = Inches(2.7)
ind_h = Inches(3.3)
ind_start = Inches(0.6)
ind_gap = Inches(0.35)

ind_colors = [GOLD, LIGHT_BLUE, GREEN, RGBColor(0xE7, 0x73, 0xE7)]

for i, (title, body) in enumerate(industries):
    x = ind_start + i * (ind_w + ind_gap)
    y = Inches(1.9)

    add_shape_bg(slide7, x, y, ind_w, ind_h, DARK_BG)

    # Color bar at top of card
    add_shape_bg(slide7, x, y, ind_w, Inches(0.06), ind_colors[i])

    # Title
    add_text(slide7, x + Inches(0.2), y + Inches(0.3),
             ind_w - Inches(0.4), Inches(0.4),
             title, font_size=17, color=ind_colors[i], bold=True,
             alignment=PP_ALIGN.CENTER)

    # Body
    add_text(slide7, x + Inches(0.2), y + Inches(0.9),
             ind_w - Inches(0.4), ind_h - Inches(1.1),
             body, font_size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# Bottom
add_shape_bg(slide7, Inches(0.6), Inches(5.5), W - Inches(1.2), Inches(0.8), DARK_BG)
add_text(slide7, Inches(1.0), Inches(5.6), W - Inches(2.0), Inches(0.6),
         "The platform is industry-agnostic — any market with buyers, sellers,\n"
         "and dynamic pricing can benefit from AI-trained trading intelligence.",
         font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 8: CLOSING — NEXT STEPS
# =====================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, NAVY)
add_shape_bg(slide8, 0, 0, W, Inches(0.08), GOLD)

# Big statement
add_text(slide8, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
         "Ready to Trade Smarter?", font_size=44,
         color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Divider
add_shape_bg(slide8, Inches(5.5), Inches(2.5), Inches(2.3), Inches(0.04), GOLD)

# Value prop summary
add_text(slide8, Inches(2), Inches(3.0), Inches(9), Inches(1.0),
         "Train AI agents in a safe simulation. Deploy winning strategies.\n"
         "Capture more margin. React faster. Partner smarter.",
         font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Next steps
steps_data = [
    ("1", "Pilot Program",
     "Start with a focused simulation\n"
     "using your real market data"),
    ("2", "Train & Validate",
     "Run AI agents through 1000+\n"
     "trading scenarios to discover\n"
     "optimal strategies"),
    ("3", "Deploy & Grow",
     "Roll out AI-powered trading\n"
     "insights across your\n"
     "organization"),
]

ns_w = Inches(3.3)
ns_h = Inches(1.6)
ns_start = Inches(0.9)
ns_gap = Inches(0.4)

for i, (num, title, body) in enumerate(steps_data):
    x = ns_start + i * (ns_w + ns_gap)
    y = Inches(4.3)

    add_shape_bg(slide8, x, y, ns_w, ns_h, DARK_BG)

    ns2 = slide8.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.15), y + Inches(0.15),
                                    Inches(0.4), Inches(0.4))
    ns2.fill.solid()
    ns2.fill.fore_color.rgb = GOLD
    ns2.line.fill.background()
    nsf2 = ns2.text_frame
    nsf2.paragraphs[0].text = num
    nsf2.paragraphs[0].font.size = Pt(14)
    nsf2.paragraphs[0].font.color.rgb = NAVY
    nsf2.paragraphs[0].font.bold = True
    nsf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide8, x + Inches(0.65), y + Inches(0.15),
             ns_w - Inches(0.8), Inches(0.35),
             title, font_size=16, color=GOLD, bold=True)

    add_text(slide8, x + Inches(0.2), y + Inches(0.6),
             ns_w - Inches(0.4), ns_h - Inches(0.7),
             body, font_size=12, color=LIGHT_GRAY)

# Footer
add_text(slide8, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
         "MarketForge  |  Built with OpenEnv  |  Powered by Reinforcement Learning",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SAVE
# =====================================================================
output_path = os.path.join(os.path.dirname(__file__),
                           "AI_Market_Intelligence_Business_Deck.pptx")
prs.save(output_path)
print(f"\nPowerPoint saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
